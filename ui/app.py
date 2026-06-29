import sys
import os
import tempfile
import shutil
from pathlib import Path
import streamlit as st

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from vera_main import ask_vera, ask_vera_stream
from ingestion.pipeline import run_ingestion
from vectorstore.retriever import get_chroma_collection

st.set_page_config(
    page_title="VERA",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap');

html, body, [class*="css"] { font-family: 'Inter', sans-serif; }
#MainMenu, footer, header { visibility: hidden; }
.stApp { background: #1E1A1B; color: #d4b8bc; }
[data-testid="stSidebar"] { background: #2A2326; border-right: 1px solid #3B2A2F; }

/* Center & Compact Main Workspace Container */
.main .block-container {
    max-width: 850px !important;
    margin: 0 auto !important;
    padding-top: 1.5rem !important;
    padding-bottom: 1.5rem !important;
}

/* Hide "Press Enter to submit" */
.stForm small, [data-testid="InputInstructions"] { display: none !important; }

.vera-logo {
    font-family: 'JetBrains Mono', monospace;
    font-size: 22px; font-weight: 500;
    color: #C8A4A8; letter-spacing: 3px; padding: 8px 0 2px 0;
}
.vera-tagline {
    font-size: 11px; color: #5C3A40; letter-spacing: 1px;
    text-transform: uppercase; padding-bottom: 16px;
    border-bottom: 1px solid #3B2A2F; margin-bottom: 16px;
}

/* Dashboard cards */
.dash-section { margin-bottom: 20px; }
.dash-title { font-size: 11px; font-weight: 600; letter-spacing: 1px; text-transform: uppercase; color: #5C3A40; margin-bottom: 10px; }
.doc-card {
    background: #2A2326; border: 1px solid #3B2A2F; border-radius: 10px;
    padding: 10px 14px; margin: 6px 0; display: flex; align-items: center; gap: 10px;
}
.doc-icon { font-size: 18px; }
.doc-info { flex: 1; }
.doc-name { font-size: 13px; color: #C8A4A8; font-weight: 500; }
.doc-meta { font-size: 11px; color: #5C3A40; margin-top: 2px; }

/* Status badge */
.status-dot { width: 8px; height: 8px; border-radius: 50%; display: inline-block; margin-right: 6px; }
.status-ready { background: #68d391; box-shadow: 0 0 6px #68d391; }
.status-empty { background: #f6ad55; box-shadow: 0 0 6px #f6ad55; }
.status-bar {
    display: flex; align-items: center; gap: 8px;
    font-size: 12px; color: #8A5A61; padding: 8px 0;
}

/* Chat messages layout */
.msg-user-wrap {
    display: flex !important;
    justify-content: flex-end !important;
    margin: 8px 0 !important;
}
.msg-user {
    background: #2A2326 !important;
    border: 1px solid #3B2A2F !important;
    border-radius: 12px 12px 2px 12px !important;
    padding: 8px 14px !important;
    max-width: 80% !important;
    font-size: 13px !important;
    line-height: 1.45 !important;
    color: #d4b8bc !important;
}

.msg-vera-wrap {
    display: flex !important;
    justify-content: flex-start !important;
    margin: 8px 0 !important;
}
.msg-vera {
    background: #2A2326 !important;
    border: 1px solid #3B2A2F !important;
    border-radius: 12px 12px 12px 2px !important;
    padding: 10px 14px !important;
    max-width: 85% !important;
    font-size: 13px !important;
    line-height: 1.5 !important;
    color: #d4b8bc !important;
}
.msg-error {
    background: #1a1020; border: 1px solid #4a2035;
    border-left: 3px solid #fc8181;
    border-radius: 4px 12px 12px 12px;
    padding: 14px 18px; margin: 8px 0; margin-right: 10%;
    color: #fc8181; font-size: 13px; line-height: 1.6;
}
.msg-thinking {
    background: #231f20; border: 1px solid #3B2A2F;
    border-left: 3px solid #5C3A40;
    border-radius: 4px 12px 12px 12px;
    padding: 14px 18px; margin: 8px 0; margin-right: 10%;
    color: #5C3A40; font-size: 13px;
}
.msg-label { font-size: 9.5px; font-weight: 600; letter-spacing: 1px; text-transform: uppercase; margin-bottom: 4px; }
.label-user { color: #C8A4A8; text-align: right; }
.label-vera { color: #C8A4A8; }
.label-error { color: #fc8181; }

/* Citations */
.citations-block { margin-top: 8px; padding-top: 8px; border-top: 1px solid #3B2A2F; }
.citations-title { font-size: 9.5px; font-weight: 600; letter-spacing: 1px; text-transform: uppercase; color: #C8A4A8; margin-bottom: 4px; }
.citation-wrap { position: relative; display: inline-block; margin: 3px 4px 3px 0; }
.citation-pill {
    display: inline-block; background: #1E1A1B; border: 1px solid #3B2A2F;
    border-radius: 12px; padding: 2px 8px; font-size: 10px;
    font-family: 'JetBrains Mono', monospace; color: #C8A4A8; cursor: default;
}
.citation-tooltip {
    display: none; position: absolute; bottom: 135%; left: 0;
    width: 320px; max-width: 80vw;
    background: #2A2326; border: 1px solid #C8A4A8; border-radius: 8px;
    padding: 10px 14px; font-size: 11px; color: #d4b8bc;
    white-space: normal; word-wrap: break-word; line-height: 1.5; text-align: left;
    z-index: 99999; box-shadow: 0 6px 20px rgba(0,0,0,0.6);
}
.citation-tooltip::after {
    content: ''; position: absolute; top: 100%; left: 20px;
    border: 6px solid transparent; border-top-color: #C8A4A8;
}
.citation-wrap:hover .citation-pill { background: #5C3A40; border-color: #C8A4A8; }
.citation-wrap:hover .citation-tooltip { display: block; }

/* Input area */
.stTextInput > div > div > input {
    background: #2A2326 !important; border: 1px solid #3B2A2F !important;
    border-radius: 10px !important; color: #d4b8bc !important;
    font-size: 14px !important; padding: 12px 16px !important;
}
.stTextInput > div > div > input:focus { border-color: #C8A4A8 !important; box-shadow: none !important; }
.stTextInput > div > div > input::placeholder { color: #5C3A40 !important; }

/* Ask button — small */
.stFormSubmitButton > button {
    background: #C8A4A8 !important; border: none !important;
    color: #1E1A1B !important; border-radius: 10px !important;
    font-weight: 600 !important; font-size: 13px !important;
    padding: 8px 16px !important; height: 44px !important;
    width: auto !important;
}
.stFormSubmitButton > button:hover { background: #a0b4ff !important; }

.stButton > button {
    background: #2A2326 !important; border: 1px solid #3B2A2F !important;
    color: #8A5A61 !important; border-radius: 8px !important;
    font-size: 12px !important;
}
.stButton > button:hover { border-color: #C8A4A8 !important; color: #C8A4A8 !important; }

hr { border-color: #3B2A2F !important; }
.stToggle { font-size: 12px !important; }
/* Tighten radio button spacing */
[data-testid="stRadio"] > div { gap: 4px !important; }
[data-testid="stRadio"] label { padding: 4px 8px !important; }

/* Scroll to bottom floating arrow */
.scroll-btn {
    position: fixed !important;
    bottom: 30px !important;
    right: 30px !important;
    width: 42px !important;
    height: 42px !important;
    border-radius: 50% !important;
    background: #3B2A2F !important;
    border: 1px solid #C8A4A8 !important;
    color: #C8A4A8 !important;
    font-size: 18px !important;
    font-weight: bold !important;
    cursor: pointer !important;
    display: flex !important;
    align-items: center !important;
    justify-content: center !important;
    z-index: 999999 !important;
    box-shadow: 0 4px 12px rgba(0,0,0,0.5) !important;
    transition: all 0.2s ease !important;
}
.scroll-btn:hover {
    background: #C8A4A8 !important;
    color: #1E1A1B !important;
    transform: scale(1.1) !important;
}
</style>
""", unsafe_allow_html=True)

# ── Session state ─────────────────────────────────────────────
# Wipe ChromaDB on every fresh app start

for key, default in [
    ("messages", []),
    ("ingested_files", []),
    ("save_session", False),
    ("pending_question", None),
    ("last_mode", "Researcher"),
    ("da_active", False),
]:
    if key not in st.session_state:
        st.session_state[key] = default

# ── Helpers ───────────────────────────────────────────────────
def get_chunk_count() -> int:
    try:
        return get_chroma_collection().count()
    except:
        return 0

def clear_vectorstore():
    from utils.config import get_settings
    cfg = get_settings()
    try:
        persist_path = Path(cfg.chroma_persist_dir).resolve()
        if persist_path.exists():
            shutil.rmtree(persist_path)
            persist_path.mkdir(parents=True, exist_ok=True)
    except:
        pass
    
if "app_initialized" not in st.session_state:
    clear_vectorstore()
    st.session_state["app_initialized"] = True


def safe_html(text: str) -> str:
    """Escape text so it never breaks HTML rendering."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def render_citations(citations):
    if not citations:
        return ""
    pills = ""
    for c in citations:
        if isinstance(c, dict):
            lbl = safe_html(c.get("label", ""))
            snip = safe_html(c.get("snippet", ""))[:280]
            if len(c.get("snippet", "")) > 280:
                snip += "..."
        else:
            lbl = safe_html(str(c))
            snip = safe_html(str(c))
        pills += f'''<span class="citation-wrap">
            <span class="citation-pill">{lbl}</span>
            <span class="citation-tooltip">📄 <strong>Source Excerpt:</strong><br>{snip}</span>
        </span>'''
    return f'<div class="citations-block"><div class="citations-title">📌 Sources — hover for exact excerpt preview</div>{pills}</div>'

def render_message(role: str, content: dict):
    if role == "user":
        st.markdown(f'<div class="msg-user-wrap"><div class="msg-user"><div class="msg-label label-user">You</div>{safe_html(content["text"])}</div></div>', unsafe_allow_html=True)
    elif role == "thinking":
        thinking_placeholder = st.empty()
        for frame in ["⏳", "⌛", "⏳", "⌛"]:
            thinking_placeholder.markdown(f'<div class="msg-vera-wrap"><div class="msg-thinking"><div class="msg-label label-vera">VERA</div>{frame} Thinking about your question...</div></div>', unsafe_allow_html=True)
            import time; time.sleep(0.3)
    elif role == "error":
        st.markdown(f'<div class="msg-vera-wrap"><div class="msg-error"><div class="msg-label label-error">VERA</div>{safe_html(content["text"])}</div></div>', unsafe_allow_html=True)
    else:
        citations_html = render_citations(content.get("citations", []))
        answer = safe_html(content.get("answer", ""))
        st.markdown(f'<div class="msg-vera-wrap"><div class="msg-vera"><div class="msg-label label-vera">VERA</div><div style="white-space:pre-wrap">{answer}</div>{citations_html}</div></div>', unsafe_allow_html=True)

# ── Sidebar ───────────────────────────────────────────────────
with st.sidebar:
    st.markdown('<div class="vera-logo">VERA</div>', unsafe_allow_html=True)
    st.markdown('<div class="vera-tagline">Verifiable Evidence &amp; Reasoning Architecture</div>', unsafe_allow_html=True)

    # Status indicator
    has_docs = len(st.session_state.ingested_files) > 0
    if has_docs:
        status_html = '<div class="status-bar"><span class="status-dot status-ready"></span>Ready to answer questions</div>'
    else:
        status_html = '<div class="status-bar"><span class="status-dot status-empty"></span>Upload a document to get started</div>'
    st.markdown(status_html, unsafe_allow_html=True)

    st.markdown("<hr>", unsafe_allow_html=True)

    # Mode toggle
    st.markdown('<div class="dash-title">Mode</div>', unsafe_allow_html=True)
    mode = st.radio(
    "mode",
    options=["Researcher", "Devil's Advocate"],
    label_visibility="collapsed",
    key="vera_mode",
    )
    # Detect mode switch
    current_mode = st.session_state.get("vera_mode", "Researcher")
    prev_mode = st.session_state.get("last_mode", "Researcher")
    if current_mode != prev_mode:
        st.session_state.last_mode = current_mode
        st.session_state.messages = []
        st.session_state.pending_question = None
        st.session_state.da_active = (current_mode == "Devil's Advocate")
        st.rerun()
    st.markdown("<hr>", unsafe_allow_html=True)

    # Loaded documents dashboard (Claude-Style Collapsible Attachments)
    if st.session_state.ingested_files:
        st.markdown('<div class="dash-title">📂 Your documents</div>', unsafe_allow_html=True)
        for f in st.session_state.ingested_files:
            ext = Path(f).suffix.lower()
            if ext == ".pdf":
                icon, type_label = "📄", "PDF Document"
            elif ext in {".pptx", ".ppt"}:
                icon, type_label = "📊", "PowerPoint"
            elif ext == ".docx":
                icon, type_label = "📘", "Word Document"
            elif ext in {".png", ".jpg", ".jpeg"}:
                icon, type_label = "📷", "Image File"
            else:
                icon, type_label = "📝", "Text Note"
            st.markdown(f'''<details style="margin: 6px 0; background: #2A2326; border: 1px solid #3B2A2F; border-radius: 10px; padding: 8px 12px; font-size: 13px; color: #C8A4A8;">
                <summary style="cursor: pointer; font-weight: 500; display: flex; align-items: center; gap: 8px;">
                    <span>{icon}</span> <span>{safe_html(f)}</span> <span style="margin-left:auto; font-size:10px; background:#3B2A2F; padding:2px 8px; border-radius:12px; color:#68d391;">Ready</span>
                </summary>
                <div style="margin-top: 8px; padding-top: 8px; border-top: 1px solid #3B2A2F; font-size: 11px; color: #5C3A40;">
                    <div><strong>Format:</strong> {type_label}</div>
                    <div><strong>Status:</strong> Processed &amp; Indexed</div>
                </div>
            </details>''', unsafe_allow_html=True)
        st.markdown("<hr>", unsafe_allow_html=True)

    # Upload
    st.markdown('<div class="dash-title">➕ Add documents</div>', unsafe_allow_html=True)
    uploaded_files = st.file_uploader(
        "Upload",
        type=["pdf", "pptx", "ppt", "txt", "md", "docx", "png", "jpg", "jpeg"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )

    if "file_bytes_map" not in st.session_state:
        st.session_state.file_bytes_map = {}

    if uploaded_files:
        new_file_added = False
        for uf in uploaded_files:
            file_data = uf.read()
            st.session_state.file_bytes_map[uf.name] = file_data
            if uf.name not in st.session_state.ingested_files:
                with st.spinner(f"Reading {uf.name}..."):
                    suffix = Path(uf.name).suffix
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(file_data)
                        tmp_path = tmp.name
                    try:
                        result = run_ingestion(tmp_path, real_filename=uf.name)
                        st.session_state.ingested_files.append(uf.name)
                        st.success(f"✅ {uf.name} added")
                        new_file_added = True
                    except Exception as e:
                        st.error(f"❌ Could not read {uf.name}")
                    finally:
                        os.unlink(tmp_path)
        
        # If in Devil's Advocate mode and new document was uploaded, refresh tutor opener dynamically
        if new_file_added and st.session_state.get("vera_mode") == "Devil's Advocate":
            st.session_state.messages = []
            st.session_state.da_active = True
            st.rerun()

    st.markdown("<hr>", unsafe_allow_html=True)

    # View toggle & Session settings
    st.markdown('<div class="dash-title">🖥 Display &amp; Export</div>', unsafe_allow_html=True)
    split_view = st.checkbox("🖥 Split-Screen Document Panel", value=st.session_state.get("split_view", False), key="split_view_toggle")

    st.session_state.save_session = False

    # Generate PDF export data
    from utils.pdf_export import generate_pdf_bytes
    try:
        pdf_bytes = generate_pdf_bytes(
            st.session_state.messages,
            st.session_state.ingested_files,
            mode=st.session_state.get('vera_mode', 'Researcher')
        )
    except Exception as e:
        pdf_bytes = b""

    st.download_button(
        label="📥 Export Session Notes (PDF)",
        data=pdf_bytes,
        file_name="VERA_Study_Session.pdf",
        mime="application/pdf",
        use_container_width=True,
    )

    st.markdown("<br>", unsafe_allow_html=True)
    col1, col2 = st.columns(2)
    with col1:
        if st.button("🗑 Clear chat", use_container_width=True):
            st.session_state.messages = []
            st.session_state.pending_question = None
            st.rerun()
    with col2:
        if st.button("🔄 New session", use_container_width=True):
            if not st.session_state.save_session:
                clear_vectorstore()
            st.session_state.messages = []
            st.session_state.ingested_files = []
            st.session_state.pending_question = None
            st.rerun()

# ── Main ──────────────────────────────────────────────────────
st.markdown('''<div style="padding:8px 0 24px 0; border-bottom: 1px solid #3B2A2F; margin-bottom: 24px; display:flex; align-items:center; justify-content:space-between;">
    <span style="font-size:13px;color:#5C3A40;">Ask anything about your documents</span>
    <span style="font-size:10px;padding:3px 12px;border-radius:20px;background:#3B2A2F;color:#C8A4A8;border:0.5px solid #5C3A40;">{}</span>
</div>'''.format(st.session_state.get("vera_mode", "Researcher")), unsafe_allow_html=True)

# Sleek Real-JS Floating Scroll Button
import streamlit.components.v1 as components
components.html("""
<script>
(function() {
    var pWin = window.parent || window;
    var pDoc = pWin.document;
    
    function veraScrollToBottom() {
        try { pWin.scrollTo({ top: 9999999, behavior: 'smooth' }); } catch(e){}
        try { pDoc.documentElement.scrollTop = 9999999; } catch(e){}
        try { pDoc.body.scrollTop = 9999999; } catch(e){}
        
        var all = pDoc.querySelectorAll('*');
        for (var i = 0; i < all.length; i++) {
            try {
                if (all[i].scrollHeight > all[i].clientHeight + 20 && all[i].clientHeight > 50) {
                    all[i].scrollTop = all[i].scrollHeight;
                }
            } catch(e){}
        }
    }

    var existingBtn = pDoc.getElementById('vera-real-floating-btn');
    if (!existingBtn) {
        var btn = pDoc.createElement('button');
        btn.id = 'vera-real-floating-btn';
        btn.innerHTML = '↓';
        btn.title = 'Scroll to bottom';
        btn.setAttribute('style', 'position:fixed !important; bottom:35px !important; right:35px !important; width:44px !important; height:44px !important; border-radius:50% !important; background:#3B2A2F !important; border:1.5px solid #C8A4A8 !important; color:#C8A4A8 !important; font-size:20px !important; font-weight:bold !important; cursor:pointer !important; display:flex !important; align-items:center !important; justify-content:center !important; z-index:9999999 !important; box-shadow:0 4px 16px rgba(0,0,0,0.6) !important; transition:all 0.2s ease !important;');
        
        btn.addEventListener('mouseenter', function() {
            btn.style.background = '#C8A4A8';
            btn.style.color = '#1E1A1B';
            btn.style.transform = 'scale(1.12)';
        });
        btn.addEventListener('mouseleave', function() {
            btn.style.background = '#3B2A2F';
            btn.style.color = '#C8A4A8';
            btn.style.transform = 'scale(1)';
        });
        
        btn.addEventListener('click', function(e) {
            e.preventDefault();
            veraScrollToBottom();
        });
        
        pDoc.body.appendChild(btn);
    } else {
        existingBtn.onclick = veraScrollToBottom;
    }
})();
</script>
""", height=0)



# Split Panel Container
if st.session_state.get("split_view_toggle"):
    chat_col, doc_panel_col = st.columns([1.1, 0.9])
else:
    chat_col = st.container()
    doc_panel_col = None

# Render Document Panel FIRST so it is 100% permanently pinned on screen
if doc_panel_col is not None:
    with doc_panel_col:
        st.markdown('<div class="dash-title">📖 Document Viewer Panel</div>', unsafe_allow_html=True)
        if st.session_state.ingested_files:
            selected_file = st.selectbox("Select document to view:", st.session_state.ingested_files, key="split_doc_select")
            
            ext = Path(selected_file).suffix.lower()
            file_bytes = st.session_state.get("file_bytes_map", {}).get(selected_file)
            
            if ext in {".pdf", ".pptx", ".ppt", ".docx"}:
                st.markdown('''<div style="margin-top:6px; margin-bottom:12px; padding:8px 12px; background:#1E1A1B; border:1px solid #3B2A2F; border-radius:8px; font-size:11px; color:#C8A4A8; line-height:1.4;">
                    📌 <strong>Study Tip:</strong> To view original PowerPoint designs, Word formatting, or PDF pages, open the file on your computer alongside VERA! VERA reads all the text here to answer your questions.
                </div>''', unsafe_allow_html=True)

            if ext in {".png", ".jpg", ".jpeg"} and file_bytes:
                st.image(file_bytes, use_column_width=True)
            elif ext == ".pdf" and file_bytes:
                import base64
                b64_pdf = base64.b64encode(file_bytes).decode("utf-8")
                pdf_display = f'<iframe src="data:application/pdf;base64,{b64_pdf}" width="100%" height="600px" type="application/pdf" style="border:1px solid #3B2A2F; border-radius:8px;"></iframe>'
                st.markdown(pdf_display, unsafe_allow_html=True)
            elif ext in {".txt", ".md"} and file_bytes:
                try:
                    text_content = file_bytes.decode("utf-8", errors="ignore")
                    st.markdown(f'''<div style="background:#1E1A1B; border:1px solid #3B2A2F; border-radius:10px; padding:16px; font-size:12px; color:#d4b8bc; line-height:1.6; max-height:60vh; overflow-y:auto; white-space:pre-wrap;">
{safe_html(text_content)}
</div>''', unsafe_allow_html=True)
                except Exception:
                    st.write("Unable to display text file.")
            elif ext == ".docx" and file_bytes:
                try:
                    import docx, io
                    doc = docx.Document(io.BytesIO(file_bytes))
                    doc_text = "\n\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
                    st.markdown(f'''<div style="background:#1E1A1B; border:1px solid #3B2A2F; border-radius:10px; padding:16px; font-size:12px; color:#d4b8bc; line-height:1.6; max-height:50vh; overflow-y:auto; white-space:pre-wrap;">
{safe_html(doc_text)}
</div>''', unsafe_allow_html=True)
                except Exception as e:
                    st.error(f"Could not render Word document: {e}")
            elif ext in {".pptx", ".ppt"} and file_bytes:
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(file_bytes))
                    slides_html = []
                    for s_idx, slide in enumerate(prs.slides, 1):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        if slide_text:
                            st_content = "<br>".join([safe_html(t) for t in slide_text])
                            slides_html.append(f'<div style="background:#1E1A1B; border:1px solid #3B2A2F; border-radius:8px; padding:12px; margin-bottom:10px;"><div style="font-size:11px; color:#C8A4A8; font-weight:600; margin-bottom:6px;">📊 Slide {s_idx}</div><div style="font-size:12px; color:#d4b8bc; line-height:1.5;">{st_content}</div></div>')
                    if slides_html:
                        st.markdown('<div style="max-height:50vh; overflow-y:auto; padding-right:6px;">' + "".join(slides_html) + '</div>', unsafe_allow_html=True)
                    else:
                        st.info("PowerPoint presentation loaded (no text found on slides).")
                except Exception as e:
                    st.error(f"Could not render PowerPoint presentation: {e}")
            else:
                st.info("Uploaded material is loaded and ready.")
        else:
            st.info("No documents uploaded yet. Upload a document in the sidebar to preview it here side-by-side!")
        st.markdown('</div>', unsafe_allow_html=True)

with chat_col:
    # Welcome / empty state
    if not st.session_state.messages and not st.session_state.pending_question:
        if not has_docs:
            st.markdown('''<div style="text-align:center;padding:40px 20px;">
                <div style="font-size:36px;margin-bottom:12px;">📂</div>
                <div style="font-size:16px;color:#C8A4A8;font-weight:600;margin-bottom:6px;">No documents uploaded yet</div>
                <div style="font-size:12px;color:#d4b8bc;">Upload a document in the sidebar to get started.<br>VERA will read it and answer your questions.</div>
            </div>''', unsafe_allow_html=True)
        else:
            st.markdown('''<div style="text-align:center;padding:40px 20px;">
                <div style="font-size:36px;margin-bottom:12px;">💬</div>
                <div style="font-size:16px;color:#C8A4A8;font-weight:600;margin-bottom:6px;">Ready for your questions</div>
                <div style="font-size:12px;color:#d4b8bc;">Type a question below about your uploaded documents.</div>
            </div>''', unsafe_allow_html=True)

    # Devil's Advocate opening message
    if st.session_state.get("da_active") and not st.session_state.messages:
        docs = st.session_state.ingested_files
        if docs:
            doc_list = ", ".join(docs)
            opener = f"I've read your documents: **{doc_list}**. What topic would you like to be tested on? You can name a specific concept, or I can pick one for you."
        else:
            opener = "No documents are uploaded yet. Please upload a PDF or PowerPoint first, then switch to Devil's Advocate mode."
        st.session_state.messages.append({
            "role": "vera",
            "content": {
                "answer": opener,
                "citations": [],
                "chunks_used": 0,
            }
        })
        st.session_state.da_active = False
        st.rerun()

    # Render chat history
    for msg in st.session_state.messages:
        render_message(msg["role"], msg["content"])

    # Handle pending question
    if st.session_state.pending_question:
        question = st.session_state.pending_question

        render_message("user", {"text": question})
        render_message("thinking", {})

        if get_chunk_count() == 0:
            st.session_state.messages.append({"role": "user", "content": {"text": question}})
            st.session_state.messages.append({"role": "error", "content": {
                "text": "No documents have been uploaded yet. Please upload a PDF or PowerPoint file using the sidebar before asking questions."
            }})
        else:
            try:
                stream_box = st.empty()
                full_answer = ""
                citations = []
                chunks_used = 0

                for event in ask_vera_stream(
                    question,
                    mode=st.session_state.get("vera_mode", "Researcher"),
                    chat_history=st.session_state.messages,
                ):
                    if event["type"] == "meta":
                        citations = event.get("citations", [])
                        chunks_used = event.get("chunks_used", 0)
                    elif event["type"] == "token":
                        full_answer += event.get("content", "")
                        citations_html = render_citations(citations)
                        stream_box.markdown(f'<div class="msg-vera-wrap"><div class="msg-vera"><div class="msg-label label-vera">VERA</div><div style="white-space:pre-wrap">{safe_html(full_answer)}</div>{citations_html}</div></div>', unsafe_allow_html=True)

                st.session_state.messages.append({"role": "user", "content": {"text": question}})
                st.session_state.messages.append({"role": "vera", "content": {
                    "answer": full_answer,
                    "citations": citations,
                    "chunks_used": chunks_used,
                }})
            except Exception as e:
                err = str(e)
                st.session_state.messages.append({"role": "user", "content": {"text": question}})
                st.session_state.messages.append({"role": "error", "content": {
                    "text": "VERA could not process your question. Please make sure Ollama is running in the background."
                }})

        st.session_state.pending_question = None
        st.rerun()

    # Input
    st.markdown("<div style='height:20px'></div>", unsafe_allow_html=True)
    with st.form(key="qform", clear_on_submit=True):
        col1, col2 = st.columns([6, 1])
        with col1:
            user_input = st.text_input(
                "q",
                placeholder="Ask a question about your documents...",
                label_visibility="collapsed",
            )
        with col2:
            submitted = st.form_submit_button("Ask →")

    if submitted and user_input.strip():
        st.session_state.pending_question = user_input.strip()
        st.rerun()