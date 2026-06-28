from pathlib import Path
from typing import List, Dict, Any

from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter

from utils.config import get_settings
from utils.logger import get_logger

log = get_logger(__name__)
settings = get_settings()

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt"}


def _load_pdf(path: Path, source_name: str) -> List[Dict[str, Any]]:
    loader = PyPDFLoader(str(path))
    raw_pages = loader.load()
    log.info(f"Loaded {len(raw_pages)} pages from {source_name}")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    chunks = splitter.split_documents(raw_pages)

    results = []
    for i, chunk in enumerate(chunks):
        results.append({
            "text": chunk.page_content.strip(),
            "source": source_name,
            "page": chunk.metadata.get("page", 0) + 1,
            "chunk_id": f"{Path(source_name).stem}_chunk_{i:04d}",
        })
    return results


def _load_pptx(path: Path, source_name: str) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(str(path))
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""],
    )

    results = []
    chunk_counter = 0
    source_stem = Path(source_name).stem.replace(" ", "_")

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        full_text = "\n".join(slide_text).strip()

        if not full_text:
            continue

        # Split slide text into chunks if needed
        sub_chunks = splitter.split_text(full_text)
        for sub in sub_chunks:
            if len(sub.strip()) > 20:
                results.append({
                    "text": sub.strip(),
                    "source": source_name,
                    "page": slide_num,  # slide number acts as page
                    "chunk_id": f"{source_stem}_chunk_{chunk_counter:04d}",
                })
                chunk_counter += 1

    log.info(f"Loaded {len(prs.slides)} slides → {len(results)} chunks from {source_name}")
    return results


def load_and_chunk_pdf(pdf_path: str, real_filename: str = None) -> List[Dict[str, Any]]:
    """
    Loads a PDF or PPTX and splits into chunks with citation metadata.

    Args:
        pdf_path: Path to the file (may be a temp path).
        real_filename: Original filename to store in citations.

    Returns:
        List of chunk dicts with text, source, page, chunk_id.
    """
    path = Path(pdf_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {pdf_path}")

    source_name = real_filename or path.name
    ext = Path(source_name).suffix.lower()

    log.info(f"Loading file: {source_name} (type: {ext})")

    if ext == ".pdf":
        results = _load_pdf(path, source_name)
    elif ext in {".pptx", ".ppt"}:
        results = _load_pptx(path, source_name)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    # Filter empty chunks
    results = [r for r in results if len(r["text"]) > 20]
    log.info(f"Final chunk count after filtering: {len(results)}")
    return results
